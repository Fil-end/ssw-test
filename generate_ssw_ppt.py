import sys
import os
os.add_dll_directory(r'D:\ProgramData\anaconda3\envs\Filend\bin')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


class SSWPPTGenerator:
    def __init__(self, output_path="ssw_method_presentation.pptx"):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.output_path = output_path

        self.primary = RGBColor(0x1B, 0x3A, 0x5C)
        self.secondary = RGBColor(0x2E, 0x75, 0xB6)
        self.accent = RGBColor(0x00, 0x70, 0xC0)
        self.light_blue = RGBColor(0xD6, 0xE4, 0xF0)
        self.lighter_blue = RGBColor(0xEB, 0xF3, 0xFA)
        self.text_dark = RGBColor(0x33, 0x33, 0x33)
        self.text_white = RGBColor(0xFF, 0xFF, 0xFF)
        self.bg_white = RGBColor(0xFF, 0xFF, 0xFF)

        self.slide_counter = 0

    def _add_background_white(self, slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.bg_white

    def _add_accent_bar(self, slide, left=0, top=0, width=Inches(0.12), height=None, color=None):
        if height is None:
            height = self.prs.slide_height
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color or self.accent
        shape.line.fill.background()

    def _add_header_bar(self, slide):
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                        self.prs.slide_width, Inches(1.1))
        header.fill.solid()
        header.fill.fore_color.rgb = self.primary
        header.line.fill.background()

        bottom_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.1),
                                             self.prs.slide_width, Inches(0.04))
        bottom_line.fill.solid()
        bottom_line.fill.fore_color.rgb = self.accent
        bottom_line.line.fill.background()

    def _add_text_box(self, slide, left, top, width, height, text,
                      font_size=18, bold=False, color=None, alignment=PP_ALIGN.LEFT,
                      font_name="Microsoft YaHei"):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color or self.text_dark
        p.font.name = font_name
        p.alignment = alignment
        return txBox

    def _add_bullet_list(self, slide, left, top, width, height, items,
                         font_size=16, color=None, spacing=Pt(10)):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            if isinstance(item, dict):
                p.text = ""
                run = p.add_run()
                run.text = "\u2022 "
                run.font.size = Pt(font_size)
                run.font.color.rgb = self.accent
                run.font.name = "Microsoft YaHei"

                run2 = p.add_run()
                run2.text = item["text"]
                run2.font.size = Pt(font_size)
                run2.font.bold = item.get("bold", False)
                run2.font.color.rgb = item.get("color", color or self.text_dark)
                run2.font.name = "Microsoft YaHei"
            else:
                p.text = "\u2022 " + item
                p.font.size = Pt(font_size)
                p.font.color.rgb = color or self.text_dark
                p.font.name = "Microsoft YaHei"

            p.space_after = spacing
            p.level = 0

        return txBox

    def _add_slide_number(self, slide):
        self._add_text_box(slide, Inches(12.2), Inches(0.3), Inches(0.9), Inches(0.5),
                          f"{self.slide_counter + 1}", font_size=13,
                          color=self.light_blue, alignment=PP_ALIGN.RIGHT)

    def add_title_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_background_white(slide)

        self._add_accent_bar(slide, left=Inches(0), top=Inches(0),
                            width=Inches(0.2), height=self.prs.slide_height)

        left_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                         Inches(5.5), self.prs.slide_height)
        left_bg.fill.solid()
        left_bg.fill.fore_color.rgb = self.primary
        left_bg.line.fill.background()

        self._add_text_box(slide, Inches(0.8), Inches(1.5), Inches(4.5), Inches(2),
                          "Stochastic Surface Walking", font_size=36, bold=True,
                          color=self.text_white)
        self._add_text_box(slide, Inches(0.8), Inches(3.5), Inches(4.5), Inches(1),
                          "SSW \u65b9\u6cd5\u8be6\u89e3", font_size=28, bold=True,
                          color=self.light_blue)

        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4.6),
                                      Inches(2.5), Inches(0.05))
        line.fill.solid()
        line.fill.fore_color.rgb = self.accent
        line.line.fill.background()

        self._add_text_box(slide, Inches(0.8), Inches(5.0), Inches(4.5), Inches(0.5),
                          "\u5168\u5c40\u4f18\u5316\u4e0e\u53cd\u5e94\u91c7\u6837\u7b97\u6cd5", font_size=18,
                          color=self.light_blue)

        self._add_text_box(slide, Inches(6.5), Inches(2.0), Inches(6), Inches(1),
                          "\u65b9\u6cd5\u80cc\u666f\u4e0e\u539f\u7406", font_size=24, bold=True,
                          color=self.primary)
        self._add_text_box(slide, Inches(6.5), Inches(3.0), Inches(6), Inches(0.5),
                          "\u57fa\u4e8e MACE \u673a\u5668\u529b\u573a\u7684\u9ad8\u6548\u5b9e\u73b0", font_size=18,
                          color=self.secondary)

        self._add_text_box(slide, Inches(6.5), Inches(4.0), Inches(6), Inches(2.5),
                          "\u2022 \u968f\u673a\u8868\u9762\u884c\u8d70\u7b97\u6cd5 (SSW)\n"
                          "\u2022 \u5c40\u90e8\u8f6f\u5316 SSW (LS-SSW)\n"
                          "\u2022 Biased Dimer Rotation \u6280\u672f\n"
                          "\u2022 \u81ea\u9002\u5e94\u7f5a\u52bf\u8c03\u8282\u673a\u5236\n"
                          "\u2022 \u6279\u91cf\u5e76\u884c\u8ba1\u7b97\u52a0\u901f",
                          font_size=16, color=self.text_dark)

        self.slide_counter += 1
        return slide

    def add_outline_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_background_white(slide)
        self._add_accent_bar(slide, left=Inches(0), top=Inches(0),
                            width=Inches(0.12), height=self.prs.slide_height)
        self._add_header_bar(slide)
        self._add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.7),
                          "\u76ee\u5f55 / Outline", font_size=28, bold=True,
                          color=self.text_white)
        self._add_slide_number(slide)

        sections = [
            ("01", "SSW \u65b9\u6cd5\u80cc\u666f\u4e0e\u52a8\u673a"),
            ("02", "\u6838\u5fc3\u7b97\u6cd5\u6d41\u7a0b"),
            ("03", "Biased Dimer Rotation \u6280\u672f"),
            ("04", "\u5c40\u90e8\u8f6f\u5316 SSW (LS-SSW)"),
            ("05", "\u81ea\u9002\u5e94\u7f5a\u52bf\u8c03\u8282\u673a\u5236"),
            ("06", "MACE \u673a\u5668\u529b\u573a\u96c6\u6210"),
            ("07", "\u6279\u91cf\u8ba1\u7b97\u4e0e\u6027\u80fd\u4f18\u5316"),
            ("08", "\u8499\u7279\u5361\u6d1b\u63a5\u53d7\u51c6\u5219"),
            ("09", "\u5b9e\u73b0\u7ec6\u8282\u4e0e\u5173\u952e\u53c2\u6570"),
            ("10", "\u603b\u7ed3\u4e0e\u5c55\u671b"),
        ]

        for i, (num, title) in enumerate(sections):
            y_pos = Inches(1.5) + Inches(0.55 * i)

            num_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), y_pos,
                                           Inches(0.7), Inches(0.42))
            num_bg.fill.solid()
            num_bg.fill.fore_color.rgb = self.accent if i % 2 == 0 else self.secondary
            num_bg.line.fill.background()

            self._add_text_box(slide, Inches(1.2), y_pos + Inches(0.03), Inches(0.7), Inches(0.4),
                              num, font_size=14, bold=True,
                              color=self.text_white, alignment=PP_ALIGN.CENTER)

            self._add_text_box(slide, Inches(2.2), y_pos + Inches(0.03), Inches(9), Inches(0.4),
                              title, font_size=18, color=self.text_dark)

        self.slide_counter += 1
        return slide

    def add_content_slide(self, title, bullets, notes=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_background_white(slide)
        self._add_accent_bar(slide, left=Inches(0), top=Inches(0),
                            width=Inches(0.12), height=self.prs.slide_height)
        self._add_header_bar(slide)
        self._add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.7),
                          title, font_size=26, bold=True, color=self.text_white)
        self._add_slide_number(slide)

        self._add_bullet_list(slide, Inches(0.8), Inches(1.4), Inches(11.5), Inches(5.5),
                             bullets, font_size=17, spacing=Pt(12))

        if notes:
            note_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.3),
                                              Inches(11.5), Inches(0.7))
            note_box.fill.solid()
            note_box.fill.fore_color.rgb = self.lighter_blue
            note_box.line.color.rgb = self.secondary
            note_box.line.width = Pt(1)
            self._add_text_box(slide, Inches(1.0), Inches(6.35), Inches(11), Inches(0.6),
                              notes, font_size=13, color=self.secondary)

        self.slide_counter += 1
        return slide

    def add_two_column_slide(self, title, left_title, left_bullets, right_title, right_bullets, notes=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_background_white(slide)
        self._add_accent_bar(slide, left=Inches(0), top=Inches(0),
                            width=Inches(0.12), height=self.prs.slide_height)
        self._add_header_bar(slide)
        self._add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.7),
                          title, font_size=26, bold=True, color=self.text_white)
        self._add_slide_number(slide)

        divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.6), Inches(1.5),
                                         Inches(0.03), Inches(5))
        divider.fill.solid()
        divider.fill.fore_color.rgb = self.light_blue
        divider.line.fill.background()

        self._add_text_box(slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.5),
                          left_title, font_size=20, bold=True, color=self.primary)
        self._add_bullet_list(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5),
                             left_bullets, font_size=15, spacing=Pt(10))

        self._add_text_box(slide, Inches(7.0), Inches(1.4), Inches(5.5), Inches(0.5),
                          right_title, font_size=20, bold=True, color=self.primary)
        self._add_bullet_list(slide, Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.5),
                             right_bullets, font_size=15, spacing=Pt(10))

        if notes:
            note_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.3),
                                              Inches(11.5), Inches(0.7))
            note_box.fill.solid()
            note_box.fill.fore_color.rgb = self.lighter_blue
            note_box.line.color.rgb = self.secondary
            note_box.line.width = Pt(1)
            self._add_text_box(slide, Inches(1.0), Inches(6.35), Inches(11), Inches(0.6),
                              notes, font_size=13, color=self.secondary)

        self.slide_counter += 1
        return slide

    def add_formula_slide(self, title, formulas, bullets=None, notes=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_background_white(slide)
        self._add_accent_bar(slide, left=Inches(0), top=Inches(0),
                            width=Inches(0.12), height=self.prs.slide_height)
        self._add_header_bar(slide)
        self._add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.7),
                          title, font_size=26, bold=True, color=self.text_white)
        self._add_slide_number(slide)

        for i, formula in enumerate(formulas):
            y_pos = Inches(1.5) + Inches(1.2 * i)

            formula_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), y_pos,
                                                 Inches(11), Inches(0.9))
            formula_box.fill.solid()
            formula_box.fill.fore_color.rgb = self.lighter_blue
            formula_box.line.color.rgb = self.secondary
            formula_box.line.width = Pt(1.5)

            label = formula.get("label", "")
            eq = formula.get("equation", "")

            if label:
                self._add_text_box(slide, Inches(1.3), y_pos + Inches(0.05), Inches(2), Inches(0.35),
                                  label, font_size=14, bold=True, color=self.accent)

            # Add equation using PowerPoint's built-in equation editor
            # For PowerPoint, we need to use the equation object
            # Since python-pptx doesn't support equation objects directly,
            # we'll use a text box with proper math formatting
            eq_box = slide.shapes.add_textbox(Inches(3.3), y_pos + Inches(0.25), Inches(8), Inches(0.6))
            tf = eq_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            
            # Use proper math formatting without LaTeX syntax
            run = p.add_run()
            run.text = eq
            run.font.size = Pt(20)
            run.font.bold = True
            run.font.color.rgb = self.text_dark
            run.font.name = "Cambria Math"  # Math font for better equation rendering

        if bullets:
            bullet_y = Inches(1.5 + 1.2 * len(formulas) + 0.3)
            self._add_bullet_list(slide, Inches(0.8), bullet_y, Inches(11.5), Inches(3),
                                 bullets, font_size=15, spacing=Pt(8))

        if notes:
            note_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.3),
                                              Inches(11.5), Inches(0.7))
            note_box.fill.solid()
            note_box.fill.fore_color.rgb = self.lighter_blue
            note_box.line.color.rgb = self.secondary
            note_box.line.width = Pt(1)
            self._add_text_box(slide, Inches(1.0), Inches(6.35), Inches(11), Inches(0.6),
                              notes, font_size=13, color=self.secondary)

        self.slide_counter += 1
        return slide

    def add_flowchart_slide(self, title, steps, notes=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_background_white(slide)
        self._add_accent_bar(slide, left=Inches(0), top=Inches(0),
                            width=Inches(0.12), height=self.prs.slide_height)
        self._add_header_bar(slide)
        self._add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.7),
                          title, font_size=26, bold=True, color=self.text_white)
        self._add_slide_number(slide)

        for i, step in enumerate(steps):
            y_pos = Inches(1.4) + Inches(0.75 * i)

            num_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), y_pos,
                                               Inches(0.45), Inches(0.45))
            num_circle.fill.solid()
            num_circle.fill.fore_color.rgb = self.accent
            num_circle.line.fill.background()
            num_circle.text_frame.paragraphs[0].text = str(i + 1)
            num_circle.text_frame.paragraphs[0].font.size = Pt(14)
            num_circle.text_frame.paragraphs[0].font.color.rgb = self.text_white
            num_circle.text_frame.paragraphs[0].font.bold = True
            num_circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            is_bold = step.get("bold", False)
            self._add_text_box(slide, Inches(1.7), y_pos + Inches(0.03), Inches(10.5), Inches(0.5),
                              step["text"], font_size=16, bold=is_bold, color=self.text_dark)

            if i < len(steps) - 1:
                arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(1.15), y_pos + Inches(0.48),
                                               Inches(0.18), Inches(0.25))
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = self.light_blue
                arrow.line.fill.background()

        if notes:
            note_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.3),
                                              Inches(11.5), Inches(0.7))
            note_box.fill.solid()
            note_box.fill.fore_color.rgb = self.lighter_blue
            note_box.line.color.rgb = self.secondary
            note_box.line.width = Pt(1)
            self._add_text_box(slide, Inches(1.0), Inches(6.35), Inches(11), Inches(0.6),
                              notes, font_size=13, color=self.secondary)

        self.slide_counter += 1
        return slide

    def add_section_divider(self, section_title, section_number):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_background_white(slide)

        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                    self.prs.slide_width, self.prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.primary
        bg.line.fill.background()

        self._add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(1),
                          f"Section {section_number}", font_size=24, color=self.light_blue)
        self._add_text_box(slide, Inches(1.5), Inches(3.2), Inches(10), Inches(1.5),
                          section_title, font_size=36, bold=True, color=self.text_white)

        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(4.8),
                                      Inches(3), Inches(0.06))
        line.fill.solid()
        line.fill.fore_color.rgb = self.accent
        line.line.fill.background()

        self.slide_counter += 1
        return slide

    def add_conclusion_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_background_white(slide)
        self._add_accent_bar(slide, left=Inches(0), top=Inches(0),
                            width=Inches(0.12), height=self.prs.slide_height)
        self._add_header_bar(slide)
        self._add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.7),
                          "\u603b\u7ed3\u4e0e\u5c55\u671b", font_size=26, bold=True,
                          color=self.text_white)
        self._add_slide_number(slide)

        self._add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(0.5),
                          "\u6838\u5fc3\u8d21\u732e", font_size=20, bold=True, color=self.primary)

        self._add_bullet_list(slide, Inches(0.8), Inches(2.0), Inches(11), Inches(2),
                             [
                                 {"text": "SSW\u65b9\u6cd5\u80fd\u591f\u7cfb\u7edf\u6027\u63a2\u7d22\u52bf\u80fd\u9762\uff0c\u53d1\u73b0\u65b0\u7684\u6781\u5c0f\u70b9\u7ed3\u6784", "bold": True},
                                 {"text": "LS-SSW\u901a\u8fc7\u5c40\u90e8\u8f6f\u5316\u964d\u4f4e\u80fd\u5792\uff0c\u63d0\u9ad8\u8df3\u8dc3\u6548\u7387", "bold": True},
                                 {"text": "\u81ea\u9002\u5e94\u7f5a\u52bf\u8c03\u8282\u5b9e\u73b0\u52a8\u6001\u5e73\u8861\u63a2\u7d22\u4e0e\u5229\u7528", "bold": False},
                                 {"text": "MACE\u673a\u5668\u529b\u573a\u63d0\u4f9bDFT\u7ea7\u7cbe\u5ea6\u4e14\u8ba1\u7b97\u9ad8\u6548", "bold": False},
                                 {"text": "\u6279\u91cf\u5e76\u884c\u8ba1\u7b97\u5927\u5e45\u63d0\u5347\u6574\u4f53\u4f18\u5316\u901f\u5ea6", "bold": False},
                             ], font_size=16, spacing=Pt(10))

        self._add_text_box(slide, Inches(0.8), Inches(4.2), Inches(11), Inches(0.5),
                          "\u672a\u6765\u5c55\u671b", font_size=20, bold=True, color=self.primary)

        self._add_bullet_list(slide, Inches(0.8), Inches(4.8), Inches(11), Inches(1.5),
                             [
                                 {"text": "\u96c6\u6210\u66f4\u591a\u673a\u5668\u529b\u573a\u6a21\u578b (NequIP, M3GNet)", "bold": False},
                                 {"text": "\u652f\u6301\u66f4\u5927\u4f53\u7cfb\u548c\u66f4\u590d\u6742\u7684\u50ac\u5316\u53cd\u5e94", "bold": False},
                                 {"text": "\u7ed3\u5408\u4e3b\u52a8\u5b66\u4e60\u7b56\u7565\u81ea\u52a8\u7b5b\u9009\u5019\u9009\u7ed3\u6784", "bold": False},
                             ], font_size=16, spacing=Pt(10))

        take_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.0),
                                         Inches(11.5), Inches(0.8))
        take_box.fill.solid()
        take_box.fill.fore_color.rgb = self.lighter_blue
        take_box.line.color.rgb = self.accent
        take_box.line.width = Pt(1.5)

        self._add_text_box(slide, Inches(1.2), Inches(6.1), Inches(10.5), Inches(0.6),
                          "\u6838\u5fc3\u4ef7\u503c\uff1aSSW\u63d0\u4f9b\u4e86\u4e00\u79cd\u7cfb\u7edf\u6027\u63a2\u7d22\u52bf\u80fd\u9762\u7684\u901a\u7528\u6846\u67b6\uff0c\u7ed3\u5408ML\u52bf\u51fd\u53ef\u5b9e\u73b0\u9ad8\u6548\u5168\u5c40\u4f18\u5316",
                          font_size=15, bold=True, color=self.primary, alignment=PP_ALIGN.CENTER)

        self.slide_counter += 1
        return slide

    def add_qa_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_background_white(slide)

        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                    self.prs.slide_width, self.prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.primary
        bg.line.fill.background()

        self._add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(2),
                          "\u611f\u8c22\u8046\u542c\uff01", font_size=48, bold=True,
                          color=self.text_white, alignment=PP_ALIGN.CENTER)

        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(4.5),
                                      Inches(2.5), Inches(0.06))
        line.fill.solid()
        line.fill.fore_color.rgb = self.accent
        line.line.fill.background()

        self._add_text_box(slide, Inches(1.5), Inches(5.0), Inches(10), Inches(1.5),
                          "Questions & Discussion", font_size=24,
                          color=self.light_blue, alignment=PP_ALIGN.CENTER)

        self.slide_counter += 1
        return slide

    def save(self):
        self.prs.save(self.output_path)
        print(f"PPT\u5df2\u4fdd\u5b58\u81f3: {self.output_path}")
        print(f"\u603b\u5e7b\u706f\u7247\u6570: {self.slide_counter}")
        return self.output_path


def generate_ssw_ppt():
    os.makedirs("ssw_output", exist_ok=True)
    gen = SSWPPTGenerator(output_path="ssw_output/SSW_Method_Presentation_LaTeX.pptx")

    gen.add_title_slide()
    gen.add_outline_slide()

    gen.add_section_divider("SSW \u65b9\u6cd5\u80cc\u666f\u4e0e\u52a8\u673a", "01")

    gen.add_content_slide(
        title="\u4ec0\u4e48\u662f SSW \u65b9\u6cd5\uff1f",
        bullets=[
            {"text": "SSW (Stochastic Surface Walking) \u662f\u4e00\u79cd\u5168\u5c40\u4f18\u5316\u7b97\u6cd5\uff0c\u7528\u4e8e\u63a2\u7d22\u52bf\u80fd\u9762 (PES) \u4e0a\u7684\u6781\u5c0f\u70b9", "bold": True},
            {"text": "\u6838\u5fc3\u601d\u60f3\uff1a\u901a\u8fc7\u968f\u673a\u6270\u52a8 + \u6700\u5c0f\u80fd\u91cf\u8def\u5f84\u8ddf\u8e2a\uff0c\u7cfb\u7edf\u6027\u5730\u8df3\u8dc3\u5230\u65b0\u7684\u80fd\u91cf\u4f4e\u8c37", "bold": True},
            {"text": "\u4f20\u7edf\u65b9\u6cd5\u5c40\u9650\uff1a\u5bb9\u6613\u9677\u5165\u5c40\u90e8\u6781\u5c0f\u503c\uff0c\u65e0\u6cd5\u5168\u9762\u63a2\u7d22\u52bf\u80fd\u9762", "bold": False},
            {"text": "SSW \u4f18\u52bf\uff1a\u80fd\u591f\u8de8\u8d8a\u80fd\u5792\uff0c\u53d1\u73b0\u5168\u5c40\u6700\u4f18\u6216\u63a5\u8fd1\u5168\u5c40\u6700\u4f18\u7ed3\u6784", "bold": False},
            {"text": "\u5e94\u7528\u9886\u57df\uff1a\u50ac\u5316\u53cd\u5e94\u8def\u5f84\u641c\u7d22\u3001\u6676\u4f53\u7ed3\u6784\u9884\u6d4b\u3001\u56e2\u7c07\u7a33\u5b9a\u6784\u578b\u7b49", "bold": False},
        ],
        notes="\u5173\u952e\u70b9\uff1aSSW \u4e0d\u4f9d\u8d56\u521d\u59cb\u731c\u6d4b\uff0c\u80fd\u591f\u81ea\u52a8\u63a2\u7d22\u52bf\u80fd\u9762\u7684\u5168\u5c40\u7279\u6027"
    )

    gen.add_two_column_slide(
        title="\u4f20\u7edf\u65b9\u6cd5 vs SSW \u65b9\u6cd5",
        left_title="\u4f20\u7edf\u5c40\u90e8\u4f18\u5316",
        left_bullets=[
            {"text": "\u68af\u5ea6\u4e0b\u964d\u7c7b\u65b9\u6cd5 (LBFGS, CG)", "bold": False},
            {"text": "\u4e25\u91cd\u4f9d\u8d56\u521d\u59cb\u7ed3\u6784", "bold": True},
            {"text": "\u53ea\u80fd\u627e\u5230\u9644\u8fd1\u7684\u5c40\u90e8\u6781\u5c0f\u503c", "bold": True},
            {"text": "\u65e0\u6cd5\u8de8\u8d8a\u80fd\u5792", "bold": False},
            {"text": "\u9002\u5408\u5c40\u90e8\u7cbe\u786e\u4f18\u5316", "bold": False},
        ],
        right_title="SSW \u5168\u5c40\u4f18\u5316",
        right_bullets=[
            {"text": "\u968f\u673a\u6270\u52a8 + \u80fd\u91cf\u8def\u5f84\u8ddf\u8e2a", "bold": False},
            {"text": "\u4e0d\u4f9d\u8d56\u521d\u59cb\u7ed3\u6784", "bold": True},
            {"text": "\u80fd\u591f\u8df3\u8dc3\u5230\u65b0\u7684\u80fd\u91cf\u4f4e\u8c37", "bold": True},
            {"text": "\u901a\u8fc7\u7f5a\u52bf\u51fd\u6570\u964d\u4f4e\u80fd\u5792", "bold": False},
            {"text": "\u9002\u5408\u5168\u5c40\u7ed3\u6784\u641c\u7d22", "bold": False},
        ],
        notes="SSW \u7684\u6838\u5fc3\u4f18\u52bf\u5728\u4e8e\u80fd\u591f\u7cfb\u7edf\u6027\u5730\u63a2\u7d22\u6574\u4e2a\u52bf\u80fd\u9762\uff0c\u800c\u4e0d\u662f\u5c40\u9650\u4e8e\u5355\u4e2a\u80fd\u91cf\u4f4e\u8c37"
    )

    gen.add_section_divider("\u6838\u5fc3\u7b97\u6cd5\u6d41\u7a0b", "02")

    gen.add_flowchart_slide(
        title="SSW \u7b97\u6cd5\u6d41\u7a0b",
        steps=[
            {"text": "\u521d\u59cb\u5316\uff1a\u4ece\u5f53\u524d\u7ed3\u6784\u51fa\u53d1\uff0c\u8fdb\u884c\u5c40\u90e8\u677e\u5f1b (Local Relaxation)", "bold": True},
            {"text": "\u751f\u6210\u968f\u673a\u6a21\u5f0f\uff1a\u7ec4\u5408\u5168\u5c40\u968f\u673a\u6a21\u5f0f + \u5c40\u90e8\u53cc\u539f\u5b50\u6a21\u5f0f", "bold": True},
            {"text": "Biased Dimer Rotation\uff1a\u627e\u5230\u6700\u4f73\u8f6f\u5316\u65b9\u5411 N", "bold": True},
            {"text": "\u9ad8\u65af\u578b\u6b65\u8fdb\uff1a\u6cbf\u65b9\u5411 N \u9010\u6b65\u79fb\u52a8 (ds \u6b65\u957f)\uff0c\u66f4\u65b0\u65b9\u5411", "bold": False},
            {"text": "\u80fd\u91cf\u5224\u636e\uff1a\u5f53\u80fd\u91cf\u964d\u4f4e\u8d85\u8fc7\u9608\u503c\u65f6\u505c\u6b62\u6b65\u8fdb", "bold": False},
            {"text": "\u5168\u5c40\u677e\u5f1b\uff1a\u5728\u771f\u5b9e PES \u4e0a\u8fdb\u884c LBFGS \u4f18\u5316\uff0c\u627e\u5230\u65b0\u6781\u5c0f\u70b9", "bold": True},
            {"text": "Monte Carlo \u63a5\u53d7\uff1a\u6839\u636e Metropolis \u51c6\u5219\u51b3\u5b9a\u662f\u5426\u63a5\u53d7\u65b0\u7ed3\u6784", "bold": True},
            {"text": "\u8fed\u4ee3\uff1a\u91cd\u590d\u4e0a\u8ff0\u6b65\u9aa4\uff0c\u76f4\u5230\u8fbe\u5230\u6700\u5927\u6b65\u6570\u6216\u6536\u655b", "bold": False},
        ],
        notes="\u6bcf\u4e00\u6b21\u5faa\u73af\u90fd\u662f\u4e00\u6b21\u201c\u63a2\u7d22-\u53d1\u73b0-\u9a8c\u8bc1\u201d\u7684\u8fc7\u7a0b"
    )

    gen.add_section_divider("Biased Dimer Rotation \u6280\u672f", "03")

    gen.add_content_slide(
        title="Biased Dimer Rotation \u539f\u7406",
        bullets=[
            {"text": "\u76ee\u6807\uff1a\u627e\u5230\u80fd\u591f\u6700\u5927\u7a0b\u5ea6\u964d\u4f4e\u80fd\u91cf\u7684\u641c\u7d22\u65b9\u5411 N", "bold": True},
            {"text": "\u521d\u59cb\u65b9\u5411\uff1a\u968f\u673a\u751f\u6210\u5168\u5c40\u6a21\u5f0f + \u5c40\u90e8\u53cc\u539f\u5b50\u6a21\u5f0f\u7684\u7ebf\u6027\u7ec4\u5408", "bold": False},
            {"text": "\u529b\u5dee\u5206\u6790\uff1a\u8ba1\u7b97\u521d\u59cb\u4f4d\u7f6e R0 \u548c\u504f\u79fb\u4f4d\u7f6e R1 \u7684\u529b\u5dee F0 - F1", "bold": True},
            {"text": "\u5782\u76f4\u5206\u91cf\uff1a\u63d0\u53d6\u529b\u5dee\u5728\u5782\u76f4\u4e8e N \u65b9\u5411\u7684\u5206\u91cf perp_F", "bold": True},
            {"text": "\u504f\u7f6e\u529b\uff1a\u5f15\u5165\u504f\u7f6e\u529b bias_force = -a * (R1-R0)N0 * N0\uff0c\u9632\u6b62\u65b9\u5411\u6f02\u79fb", "bold": True},
            {"text": "\u65b9\u5411\u66f4\u65b0\uff1aN_new = normalize(perp_F + bias_force)\uff0c\u5e73\u6ed1\u8fc7\u6e21", "bold": False},
            {"text": "\u6279\u91cf\u8ba1\u7b97\uff1a\u9884\u5148\u51c6\u5907\u6240\u6709\u7ed3\u6784\uff0c\u4e00\u6b21\u6027\u6279\u91cf\u8ba1\u7b97\u80fd\u91cf\u548c\u529b", "bold": False},
        ],
        notes="\u5173\u952e\u521b\u65b0\uff1a\u901a\u8fc7\u6279\u91cf\u8ba1\u7b97\u5927\u5e45\u52a0\u901f dimer rotation \u8fc7\u7a0b"
    )

    gen.add_formula_slide(
        title="Biased Dimer Rotation \u5173\u952e\u516c\u5f0f",
        formulas=[
            {"label": "\u529b\u5dee\u5782\u76f4\u5206\u91cf", "equation": "⊥F = (F₀ - F₁) - [(F₀ - F₁)N]N"},
            {"label": "\u504f\u7f6e\u529b", "equation": "bias_force = -a · [(R₁ - R₀)N₀] · N₀"},
            {"label": "\u65b9\u5411\u66f4\u65b0", "equation": "N_new = normalize(⊥F + bias_force)"},
            {"label": "\u5e73\u6ed1\u8fc7\u6e21", "equation": "N = normalize(0.5 · N + 0.5 · N_new)"},
        ],
        bullets=[
            {"text": "F₀: \u521d\u59cb\u4f4d\u7f6e\u7684\u529b\uff1bF₁: \u6cbf N \u65b9\u5411\u504f\u79fb δR \u540e\u7684\u529b", "bold": False},
            {"text": "a = 100.0: \u504f\u7f6e\u529b\u5f3a\u5ea6\u7cfb\u6570\uff0c\u63a7\u5236\u65b9\u5411\u7ea6\u675f\u7a0b\u5ea6", "bold": False},
            {"text": "δR = 0.01 Å: \u6253\u683c\u70b9\u95f4\u8ddd\uff0c\u7528\u4e8e\u8ba1\u7b97\u529b\u5dee", "bold": False},
        ],
        notes="\u6536\u655b\u6761\u4ef6\uff1a|⊥F| < 10⁻⁶ \u6216\u8005 |⊥F| < 10⁻⁴ \u4e14\u8fed\u4ee3\u6b21\u6570 > 3"
    )

    gen.add_section_divider("\u5c40\u90e8\u8f6f\u5316 SSW (LS-SSW)", "04")

    gen.add_content_slide(
        title="\u4ec0\u4e48\u662f LS-SSW\uff1f",
        bullets=[
            {"text": "LS-SSW = Local Softening SSW\uff0c\u5728\u6807\u51c6 SSW \u57fa\u7840\u4e0a\u5f15\u5165\u5c40\u90e8\u8f6f\u5316\u673a\u5236", "bold": True},
            {"text": "\u6838\u5fc3\u601d\u60f3\uff1a\u901a\u8fc7\u6dfb\u52a0\u6392\u65a5\u7f5a\u52bf\uff0c\u4eba\u4e3a\u201c\u8f6f\u5316\u201d\u539f\u5b50\u95f4\u7684\u76f8\u4e92\u4f5c\u7528", "bold": True},
            {"text": "\u6548\u679c\uff1a\u964d\u4f4e\u52bf\u80fd\u9762\u4e0a\u7684\u80fd\u5792\u9ad8\u5ea6\uff0c\u4f7f\u8de8\u8d8a\u80fd\u5792\u66f4\u5bb9\u6613", "bold": True},
            {"text": "\u7f5a\u52bf\u5f62\u5f0f\uff1aBuckingham \u578b\u6392\u65a5\u52bf\uff0c\u4ec5\u5bf9\u8ddd\u79bb < 3.0 \u00c5 \u7684\u539f\u5b50\u5bf9\u4f5c\u7528", "bold": False},
            {"text": "\u81ea\u9002\u5e94\u8c03\u8282\uff1a\u7f5a\u52bf\u5f3a\u5ea6 A_pq \u6839\u636e\u5f53\u524d\u7f5a\u52bf\u80fd\u91cf\u4e0e\u76ee\u6807\u503c\u7684\u5dee\u5f02\u52a8\u6001\u8c03\u6574", "bold": True},
            {"text": "\u9002\u7528\u573a\u666f\uff1a\u5168\u5c40\u4f18\u5316\u3001\u53cd\u5e94\u8def\u5f84\u91c7\u6837\u3001\u76f8\u53d8\u7814\u7a76", "bold": False},
        ],
        notes="LS-SSW \u7684\u5173\u952e\u5728\u4e8e\u5e73\u8861\u63a2\u7d22\uff08\u8f6f\u5316\u964d\u4f4e\u80fd\u5792\uff09\u548c\u5229\u7528\uff08\u771f\u5b9e PES \u4f18\u5316\uff09"
    )

    gen.add_formula_slide(
        title="Buckingham \u7f5a\u52bf\u516c\u5f0f",
        formulas=[
            {"label": "\u7f5a\u52bf\u80fd\u91cf", "equation": "Vₚ(i,j) = 0.5 · Aₚq · exp[-(r - r₀) / (ξ · r₀)]"},
            {"label": "\u7f5a\u52bf\u529b", "equation": "Fₚ(i) = (Aₚq / (ξ · r₀)) · exp[-(r - r₀) / (ξ · r₀)] · (rᵢⱼ / r)"},
            {"label": "\u81ea\u9002\u5e94\u7f29\u653e", "equation": "scale = 1.0 - λₛₜₑₚ · (Eₚₑₙₐₗₜy/atom - Yₜₐᵣgₑₜ) / 10.0"},
            {"label": "\u5f3a\u5ea6\u66f4\u65b0", "equation": "Aₚq(new) = Aₚq(old) · clip(scale, 0.5, 2.0)"},
        ],
        bullets=[
            {"text": "r₀: \u53c2\u8003\u8ddd\u79bb\uff1bξ = 0.2: \u8f6f\u5316\u5bbd\u5ea6\u53c2\u6570\uff1bλₛₜₑₚ = 1.8: \u81ea\u9002\u5e94\u6b65\u957f", "bold": False},
            {"text": "Yₜₐᵣgₑₜ = 0.02~0.1 eV/atom: \u76ee\u6807\u5e73\u5747\u7f5a\u52bf\u80fd\u91cf\uff08\u5168\u5c40\u4f18\u5316\u7528\u5c0f\u503c\uff09", "bold": False},
            {"text": "\u4ec5\u5bf9 r < 3.0 Å \u7684\u539f\u5b50\u5bf9\u4f5c\u7528\uff0c\u907f\u514d\u8fdc\u7a0b\u5e72\u6270", "bold": False},
        ],
        notes="\u7f5a\u52bf\u529b\u65b9\u5411\u4e0e\u539f\u5b50\u95f4\u6392\u65a5\u65b9\u5411\u76f8\u540c\uff0c\u8d77\u5230\u201c\u8f6f\u5316\u201d\u539f\u5b50\u95f4\u76f8\u4e92\u4f5c\u7528\u7684\u6548\u679c"
    )

    gen.add_section_divider("\u81ea\u9002\u5e94\u7f5a\u52bf\u8c03\u8282\u673a\u5236", "05")

    gen.add_content_slide(
        title="\u81ea\u9002\u5e94\u7f5a\u52bf\u8c03\u8282\u539f\u7406",
        bullets=[
            {"text": "\u95ee\u9898\uff1a\u56fa\u5b9a\u7f5a\u52bf\u5f3a\u5ea6\u65e0\u6cd5\u9002\u5e94\u4e0d\u540c\u4f53\u7cfb\u548c\u4e0d\u540c\u9632\u6bb5\u7684\u9700\u6c42", "bold": False},
            {"text": "\u89e3\u51b3\u65b9\u6848\uff1a\u6839\u636e\u5f53\u524d\u7f5a\u52bf\u80fd\u91cf\u4e0e\u76ee\u6807\u503c Y_target \u7684\u5dee\u5f02\u81ea\u52a8\u8c03\u6574", "bold": True},
            {"text": "\u5f53 E_penalty/atom > Y_target \u65f6\uff1a\u7f29\u5c0f\u7f5a\u52bf\u5f3a\u5ea6\uff0c\u907f\u514d\u8fc7\u5ea6\u6270\u52a8", "bold": True},
            {"text": "\u5f53 E_penalty/atom < Y_target \u65f6\uff1a\u589e\u5f3a\u7f5a\u52bf\u5f3a\u5ea6\uff0c\u52a0\u5feb\u63a2\u7d22\u901f\u5ea6", "bold": True},
            {"text": "\u7f5a\u52bf\u80fd\u91cf\u7531\u7f5a\u52bf\u529b\u5bf57\u6b21\u7c7b\u4f3c\u8ddd\u79bb\u6269\u5c55\u5f0f\u8ba1\u7b97", "bold": False},
            {"text": "\u8ab2\u5ea6\u6761\u4ef6\uff1aself-consistent \u7684\u52a8\u6001\u8c03\u6574\uff0c\u907f\u514d\u56fa\u5b9a\u7f5a\u52bf\u7684\u504f\u5dee", "bold": False},
        ],
        notes="\u81ea\u9002\u5e94\u673a\u5236\u4f7f LS-SSW \u80fd\u591f\u81ea\u52a8\u9002\u5e94\u4e0d\u540c\u4f53\u7cfb\u7684\u7f5a\u52bf\u5f3a\u5ea6\uff0c\u5b9e\u73b0\u66f4\u6709\u6548\u7684\u5168\u5c40\u63a2\u7d22"
    )

    gen.add_section_divider("MACE \u673a\u5668\u529b\u573a\u96c6\u6210", "06")

    gen.add_content_slide(
        title="MACE \u673a\u5668\u529b\u573a\u63d0\u4f9b\u9ad8\u6548\u8ba1\u7b97",
        bullets=[
            {"text": "MACE (Multiscale Atomic Cluster Expansion) \u662f\u4e00\u79cd\u9ad8\u7cbe\u5ea6\u7684\u673a\u5668\u529b\u573a\u6a21\u578b", "bold": True},
            {"text": "\u652f\u6301 CUDA GPU \u52a0\u901f\uff0c\u8ba1\u7b97\u901f\u5ea6\u6bd4传统 DFT \u5feb 100-1000 \u500d", "bold": True},
            {"text": "\u7ed3\u5408 DFT-D3 \u6df7\u5408\u8ba1\u7b97\uff1aMACE + TorchDFTD3 \u63d0\u4f9b\u7075\u6d3b\u7684\u51b3\u7ed3\u529f\u80fd", "bold": False},
            {"text": "\u652f\u6301\u9ad8\u7b49\u7c7b\u52a8\u7269\u8d28\uff1aFe-C \u7cfb\u7edf\u7b49\u8fc7\u6e21\u91d1\u5c5e\u7c7b\u7cfb\u7edf", "bold": False},
            {"text": "\u5df2\u96c6\u6210 BatchRelaxer\uff1a\u6279\u91cf\u5e76\u884c\u8ba1\u7b97\u591a\u4e2a\u7ed3\u6784\u7684\u80fd\u91cf\u548c\u529b", "bold": True},
            {"text": "\u652f\u6301\u6279\u91cf\u8ba1\u7b97\u6a21\u5f0f batch_calculate_properties\uff0c\u5927\u5e45\u5ea6\u52a0\u901f\u7ed3\u6784\u4f18\u5316", "bold": True},
        ],
        notes="MACE + D3 \u6df7\u5408\u8ba1\u7b97\u63d0\u4f9b\u4e86 DFT \u7ea7\u522b\u7684\u7cbe\u5ea6\uff0c\u540c\u65f6\u5177\u5907 GPU \u52a0\u901f\u7684\u9ad8\u6548\u6027"
    )

    gen.add_section_divider("\u8499\u7279\u5361\u6d1b\u63a5\u53d7\u51c6\u5219", "08")

    gen.add_content_slide(
        title="\u8499\u7279\u5361\u6d1b\u63a5\u53d7\u51c6\u5219 (Metropolis Monte Carlo)",
        bullets=[
            {"text": "\u95ee\u9898\uff1a\u5982\u679c\u65b0\u7ed3\u6784\u80fd\u91cf\u9ad8\u4e8e\u5f53\u524d\uff0c\u662f\u5426\u8de8\u8d8a\u5230\u66f4\u4f4e\u7684\u80fd\u91cf\u4f4e\u8c37\uff1f", "bold": False},
            {"text": "\u89e3\u51b3\u65b9\u6848\uff1aMetropolis Monte Carlo \u63a5\u53d7\u51c6\u5219", "bold": True},
            {"text": "\u5f53 Eₙₑw < Eₙᵤᵣᵣₑₙₜ \u65f6:\u59a5\u8fdb\u65b0\u7ed3\u6784\uff08\u80fd\u91cf\u964d\u4f4e\uff09", "bold": True},
            {"text": "\u5f53 Eₙₑw > Eₙᵤᵣᵣₑₙₜ \u65f6:\u4ee5\u6982\u7387 P = exp(-ΔE / (kᵦ · T)) \u63a5\u53d7", "bold": True},
            {"text": "\u76ee\u7684\uff1a\u5141\u8bb8\u7cfb\u7edf\u67e5\u770b\u66f4\u591a\u7684\u80fd\u91cf\u4f4e\u8c37\uff0c\u907f\u514d\u9677\u5165\u5c40\u90e8\u6781\u5c0f", "bold": True},
            {"text": "kᵦ = 8.617e-5 eV/K\uff1bT = mc_temp \u53ef\u8c03\uff08\u901a\u5e38 300-2000 K\uff09", "bold": False},
        ],
        notes="MC \u63a5\u53d7\u51c6\u5219\u4f7f\u5f97 SSW \u80fd\u591f\u63a2\u7d22\u66f4\u5e7e\u7684\u80fd\u91cf\u4f4e\u8c37\uff0c\u63d0\u9ad8\u5168\u5c40\u4f18\u5316\u7684\u5b8c\u6574\u6027"
    )

    gen.add_section_divider("\u5b9e\u73b0\u7ec6\u8282\u4e0e\u5173\u952e\u53c2\u6570", "09")

    gen.add_content_slide(
        title="\u5173\u952e\u53c2\u6570\u8bbe\u7f6e",
        bullets=[
            {"text": "temperature = 300 K\uff1aMonte Carlo \u6e29\u5ea6\uff0c\u63a7\u5236\u9ad8\u80fd\u91cfg structures \u7684\u63a5\u53d7\u6982\u7387", "bold": False},
            {"text": "ds = 0.1-0.6 \u00c5\uff1a\u9ad8\u65af\u578b\u6b65\u957f\uff0c\u6bcf\u6b21 dimer rotation \u540e\u7684\u79fb\u52a8\u8ddd\u79bb", "bold": False},
            {"text": "max_gaussians = 25-30\uff1a\u6700\u5927\u9ad8\u65af\u578b\u6b65\u6570\uff0c\u63a7\u5236\u5355\u6b21 climbing \u7684\u6b65\u6570", "bold": False},
            {"text": "mc_temp = 1000-2000 K\uff1aMonte Carlo \u6e29\u5ea6\uff0c\u8f83\u9ad8\u503c\u66f4\u5bb9\u6613\u63a5\u53d7\u9ad8\u80fd\u7ed3\u6784", "bold": False},
            {"text": "Y_target = 0.02-0.1 eV/atom\uff1aLS-SSW \u76ee\u6807\u7f5a\u52bf\u80fd\u91cf\uff0c\u5168\u5c40\u4f18\u5316\u7528\u5c0f\u503c", "bold": True},
            {"text": "xi = 0.2\uff1a\u8f6f\u5316\u5bbd\u5ea6\u53c2\u6570\uff0c\u63a7\u5236\u7f5a\u52bf\u7684\u4f5c\u7528\u8303\u56f4", "bold": False},
            {"text": "lambda_step = 1.8\uff1a\u81ea\u9002\u5e94\u6b65\u957f\uff0c\u63a7\u5236\u7f5a\u52bf\u5f3a\u5ea6\u7684\u8c03\u6574\u901f\u5ea6", "bold": False},
        ],
        notes="\u53c2\u6570\u8bbe\u7f6e\u9700\u8981\u6839\u636e\u5177\u4f53\u4f53\u7cfb\u8fdb\u884c\u8c03\u6574\uff0c\u5efa\u8bae\u5148\u4f7f\u7528\u9ed8\u8ba4\u503c\u8fdb\u884c\u6d4b\u8bd5"
    )

    gen.add_conclusion_slide()
    gen.add_qa_slide()

    gen.save()
    return gen.output_path


if __name__ == "__main__":
    output = generate_ssw_ppt()
    print(f"\nGenerated PPT: {output}")
